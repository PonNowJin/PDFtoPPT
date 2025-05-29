from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
import json
import os
from PIL import Image
import random

def json_to_pptx(json_data, output_path="output.pptx", output_watermark_path="output_watermark.pptx",image_dir=None, template='template/template_1.pptx',premium=False):
    # jsonlike_txt to json
    with open(json_data, "r", encoding="utf-8") as f:
        try:
            content = f.read().strip()
            if not content:
                raise ValueError("檔案為空，無法解析")
            data = json.loads(content)
            # print(data)
        except json.JSONDecodeError as e:
            print("JSON 解析錯誤：", e)
        except Exception as e:
            print("其他錯誤：", e)

    # init ppt
    prs = Presentation(template)

    for idx, slide_data in enumerate(data):
        # print(slide_data, flush=True)
        # 封面頁
        if idx == 0:
            # 新增一頁標題投影片
            slide = prs.slides.add_slide(prs.slide_layouts[0])

            # 設定標題與副標題文字
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白版面

        # 加入標題
        if "title" in slide_data:
            title_shape = slide.shapes.title
            if not title_shape:
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
            title_tf = title_shape.text_frame
            title_tf.text = slide_data["title"]
            title_tf.paragraphs[0].font.size = Pt(32)
            title_tf.paragraphs[0].font.bold = True

        # 加入內文（子彈點）
        if "body" in slide_data:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            
            font_size = 22  # pt
            line_spacing = 1.5

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True

            for idx, bullet in enumerate(slide_data["body"]):
                p = tf.add_paragraph() if idx != 0 else tf.paragraphs[0]
                p.text = bullet
                p.level = 0
                p.font.size = Pt(font_size)
                
            # 根據文字內容估算高度
            text_height = estimate_textbox_height(
                slide_data["body"],
                font_size_pt=22,
                line_spacing=line_spacing,
                textbox_width_inch=8  # 你設定的寬度
            )
            # print(text_height)

        # 加入圖片
        if "image" in slide_data and image_dir:
            image_path = os.path.join(image_dir, os.path.basename(slide_data["image"]))
            if os.path.exists(image_path):
                # 取得簡報寬度與高度（單位：EMU）
                slide_width_emu = prs.slide_width
                slide_height_emu = prs.slide_height

                # 換算成英吋
                slide_width_inch = slide_width_emu / Inches(1)
                slide_height_inch = slide_height_emu / Inches(1)
                
                # 最大圖片區域高度
                total_slide_height = Inches(slide_height_inch)  # 你可根據 slide 高度決定
                image_top = top + text_height + Inches(0.5)  # 空點間距
                max_height = total_slide_height - image_top - Inches(0.3)
                max_width = Inches(7)
                left = Inches(1)

                # 圖片原始尺寸與縮放
                with Image.open(image_path) as img:
                    img_width_px, img_height_px = img.size
                    width = Inches(img_width_px / 96)
                    height = Inches(img_height_px / 96)

                    scale = min(max_width / width, max_height / height)
                    new_width = width * scale
                    new_height = height * scale

                    left_centered = left + (max_width - new_width) / 2
                slide.shapes.add_picture(image_path, left_centered, image_top, width=new_width, height=new_height)
            '''
            if os.path.exists(image_path):

                
                # 設定最大顯示區塊
                max_width = Inches(7)
                max_height = Inches(4.5)
                top = Inches(3)
                left = Inches(1)

                # 圖片原始尺寸
                with Image.open(image_path) as img:
                    img_width_px, img_height_px = img.size
                    # PPTX 的 DPI 預設為 96
                    width = Inches(img_width_px / 96)
                    height = Inches(img_height_px / 96)

                    # 等比縮放
                    scale_w = max_width / width
                    scale_h = max_height / height
                    scale = min(scale_w, scale_h)  # 選擇較小的縮放比例以完全容納圖片

                    new_width = width * scale
                    new_height = height * scale

                    # 置中放置於框中
                    left_centered = left + (max_width - new_width) / 2
                    top_centered = top + (max_height - new_height) / 2

                # 插入圖片
                slide.shapes.add_picture(image_path, left_centered, top_centered, width=new_width, height=new_height)
                '''
                
        # 加入備忘稿文字
        if "description" in slide_data:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["description"]
    
    delete_slide(prs, 0)  # 刪除第一頁（index 0）（模板頁面）
    set_fonts_for_all_textboxes(prs)
    prs.save(output_path)
    print(f"PowerPoint 檔案儲存為：{output_path}")
    
    # 不是付費會員，加浮水印
    if not premium:
        watermark(output_path, output_watermark_path)


# 估算文字匡可能高度
def estimate_textbox_height(paragraphs, font_size_pt, line_spacing=1.2, textbox_width_inch=8):
    # 估算每行可容納的字數（根據 textbox 寬度與字體大小）
    chars_per_line = int((textbox_width_inch * 72) / font_size_pt * 1.6)

    line_height = font_size_pt * line_spacing
    total_lines = 0

    for text in paragraphs:
        text_len = len(text)
        line_count = max(1, (text_len + chars_per_line - 1) // chars_per_line)
        total_lines += line_count

    total_height_pt = total_lines * line_height
    return Inches(total_height_pt / 72)



def watermark(ppt_file, output_path,watermark_path='watermark/watermark.png', transparent_img='watermark/img_watermark'):
    # 製作半透明圖片
    def create_transparent_image(input_path, output_path, alpha=70):
        img = Image.open(input_path).convert("RGBA")
        # 修改 alpha（0 ~ 255）→ 100 表示大約 40% 不透明
        new_data = []
        for item in img.getdata():
            new_data.append((item[0], item[1], item[2], alpha))
        img.putdata(new_data)
        img.save(output_path, "PNG")

    # 路徑
    create_transparent_image(watermark_path, transparent_img)

    # 開啟簡報，加入浮水印
    prs = Presentation(ppt_file)
    
    # 浮水印尺寸（根據需求可調整）
    img_width = Inches(1.5)
    img_height = Inches(1.5)

    for slide in prs.slides:
        watermark_num = random.randint(2,3)

        for _ in range(watermark_num):
            x = Inches(random.uniform(0, 8.5))  # 預留空間避免超出邊界
            y = Inches(random.uniform(0, 6.0))

            slide.shapes.add_picture(transparent_img, x, y, width=img_width, height=img_height)

    # 儲存結果
    prs.save(output_path)
    

# 刪除投影片某頁
def delete_slide(prs, slide_index):
    slide_id_list = prs.slides._sldIdLst  # 這是所有 slide 的 XML 列表
    slides = list(prs.slides._sldIdLst)
    slide = prs.slides[slide_index]
    
    # 取得這張 slide 的 rId
    rId = None
    for rel in prs.part.rels:
        if prs.part.rels[rel].target_part == slide.part:
            rId = rel
            break

    if rId is None:
        raise ValueError("找不到 slide 的 rId")

    # 在 XML 中移除對應的 slide entry
    for sldId in slide_id_list:
        if sldId.get(qn("r:id")) == rId:
            slide_id_list.remove(sldId)
            break

    # 同時從 rels 中刪除連結（避免殘留）
    prs.part.drop_rel(rId)

        

# 設定整份簡報字體
def set_fonts_for_all_textboxes(prs, font_latin="Arial", font_east_asian="標楷體"):
    for num, slide in enumerate(prs.slides):
        if num == 0:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    rPr = run._r.get_or_add_rPr()
                    rPr.set(qn("a:ea"), font_east_asian)  # 設定 East Asian font
                    run.font.name = font_latin


if __name__ == "__main__":
    json_data = 'Gemini/generate_output.txt'
    output_path = 'output_pptx/test.pptx'
    output_watermark_path = 'output_pptx/test_watermark.pptx'
    img_dir = 'Crop_imgs_demo'
    json_to_pptx(json_data, output_path, output_watermark_path, img_dir)
    